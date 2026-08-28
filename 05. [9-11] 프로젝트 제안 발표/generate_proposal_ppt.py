import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # 16:9 와이드스크린 비율 설정
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 테마 색상 정의
    COLOR_BG_DARK = RGBColor(11, 41, 102)     # 서울대 Navy / 우리금융 Dark Blue
    COLOR_PRIMARY = RGBColor(0, 104, 182)    # 우리금융 Blue
    COLOR_ACCENT = RGBColor(242, 112, 89)     # 산뜻한 코랄/오렌지 (포인트 컬러)
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_TEXT_DARK = RGBColor(30, 41, 59)   # 거의 검은 회색 (가독성 향상)
    COLOR_TEXT_MUTED = RGBColor(100, 116, 139) # 본문 보조용 회색
    
    COLOR_CARD_RED = RGBColor(253, 244, 243)  # 연한 붉은색 카드 (위험/페인포인트)
    COLOR_CARD_BLUE = RGBColor(240, 247, 254) # 연한 파란색 카드 (해결방안)
    COLOR_CARD_GRAY = RGBColor(248, 250, 252) # 연한 회색 카드 (일반 카드)
    
    # 슬라이드 마스터 레이아웃 대신 완전히 빈 슬라이드를 레이아웃으로 사용 (인덱스 6)
    blank_layout = prs.slide_layouts[6]
    
    # ----------------------------------------------------
    # SLIDE 1: 타이틀 슬라이드 (다크 테마)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    
    # 배경색 지정 (어두운 단색 사각형으로 배경 채우기)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG_DARK
    bg.line.color.rgb = COLOR_BG_DARK
    
    # 포인트 장식 선 추가
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.8), Inches(11.333), Inches(0.08))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_PRIMARY
    accent_bar.line.color.rgb = COLOR_PRIMARY
    
    # 기관 로고/소속 텍스트
    org_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.6))
    tf_org = org_box.text_frame
    tf_org.word_wrap = True
    p_org = tf_org.paragraphs[0]
    p_org.text = "2026 하반기 서울대 AI 프론티어 마스터 과정  |  개인 프로젝트 제안 발표"
    p_org.font.size = Pt(14)
    p_org.font.bold = True
    p_org.font.color.rgb = COLOR_PRIMARY
    p_org.font.name = "맑은 고딕"
    
    # 메인 타이틀
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.0))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = "Generative AI 기반 기업 여신 심사 자동화 및\n실시간 리스크 센싱 시스템"
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE
    p_title.font.name = "맑은 고딕"
    
    # 영문 서브타이틀
    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.3), Inches(11.333), Inches(0.6))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Woori Corporate Credit Underwriting Assistant: Real-time Risk Sensing & Draft Generator (Woori Copilot)"
    p_sub.font.size = Pt(14)
    p_sub.font.bold = True
    p_sub.font.color.rgb = COLOR_ACCENT
    p_sub.font.name = "Arial"
    
    # 제안자 정보
    info_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(5.0), Inches(0.8))
    tf_info = info_box.text_frame
    p_info = tf_info.paragraphs[0]
    p_info.text = "제안기관: 우리금융그룹 / 우리은행\n제안자: 교육생"
    p_info.font.size = Pt(13)
    p_info.font.color.rgb = COLOR_WHITE
    p_info.font.name = "맑은 고딕"
    
    # ----------------------------------------------------
    # 공통 함수: 본문 슬라이드 템플릿 (헤더 구성)
    # ----------------------------------------------------
    def add_slide_header(title_text, slide_num_text):
        new_slide = prs.slides.add_slide(blank_layout)
        
        # 얇은 탑 그라데이션 대신 심플하고 세련된 미니 상단 네이비 바
        header_bar = new_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = COLOR_BG_DARK
        header_bar.line.color.rgb = COLOR_BG_DARK
        
        # 슬라이드 대분류 번호 및 메인 타이틀
        title_box = new_slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(10.0), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{slide_num_text}.  {title_text}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = "맑은 고딕"
        
        # 우측 상단 소속 표시
        sub_title_box = new_slide.shapes.add_textbox(Inches(10.5), Inches(0.35), Inches(2.2), Inches(0.4))
        p_sub = sub_title_box.text_frame.paragraphs[0]
        p_sub.text = "WOORI FINANCIAL GROUP"
        p_sub.alignment = PP_ALIGN.RIGHT
        p_sub.font.size = Pt(10)
        p_sub.font.bold = True
        p_sub.font.color.rgb = COLOR_PRIMARY
        p_sub.font.name = "Arial"
        
        return new_slide

    # ----------------------------------------------------
    # SLIDE 2: 주제 선정 배경 및 문제 정의
    # ----------------------------------------------------
    slide2 = add_slide_header("프로젝트 주제 및 문제 정의 (Topic & Pain Point)", "01")
    
    # 좌측 페인 포인트 카드 (붉은색 톤)
    card_left = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3))
    card_left.fill.solid()
    card_left.fill.fore_color.rgb = COLOR_CARD_RED
    card_left.line.color.rgb = COLOR_ACCENT
    card_left.line.width = Pt(1.5)
    
    left_title_box = slide2.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(5.2), Inches(0.5))
    p_lt = left_title_box.text_frame.paragraphs[0]
    p_lt.text = "🚨 현업의 Pain Point"
    p_lt.font.size = Pt(18)
    p_lt.font.bold = True
    p_lt.font.color.rgb = COLOR_ACCENT
    p_lt.font.name = "맑은 고딕"
    
    left_content_box = slide2.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(5.2), Inches(4.2))
    tf_lc = left_content_box.text_frame
    tf_lc.word_wrap = True
    
    bullets_lc = [
        "심사역의 과도한 업무 로드 및 생산성 한계",
        "  - 수십 장의 DART 공시, 산업동향 보고서, 부정 뉴스를 수작업 수집",
        "  - 기업당 보고서(Credit Memo) 초안 작성에만 최소 2~3일 소요",
        "정성적/비정형 리스크 실시간 모니터링 공백",
        "  - 배임, 소송, 특허 분쟁 등은 정형 재무제표에 뒤늦게 반영됨",
        "  - 실시간 리스크 감지 지연 시 거액의 여신 부실화 손실 초래",
        "설명 가능성(Explainability)이 없는 AI의 한계",
        "  - 여신 심사는 '거절 근거'를 법적으로 소상히 설명해야 함",
        "  - 블랙박스 형태의 LLM 단독 판정은 신뢰도 부족으로 실무 적용 불가"
    ]
    for i, b in enumerate(bullets_lc):
        p = tf_lc.paragraphs[0] if i == 0 else tf_lc.add_paragraph()
        p.text = b
        p.font.size = Pt(13) if b.startswith("  -") else Pt(14)
        p.font.bold = not b.startswith("  -")
        p.font.color.rgb = COLOR_TEXT_DARK
        p.font.name = "맑은 고딕"
        if b.startswith("  -"):
            p.font.color.rgb = COLOR_TEXT_MUTED

    # 우측 해결 방안 카드 (푸른색 톤)
    card_right = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.3))
    card_right.fill.solid()
    card_right.fill.fore_color.rgb = COLOR_CARD_BLUE
    card_right.line.color.rgb = COLOR_PRIMARY
    card_right.line.width = Pt(1.5)
    
    right_title_box = slide2.shapes.add_textbox(Inches(7.2), Inches(1.7), Inches(5.2), Inches(0.5))
    p_rt = right_title_box.text_frame.paragraphs[0]
    p_rt.text = "💡 해결 방안 (Our Solution)"
    p_rt.font.size = Pt(18)
    p_rt.font.bold = True
    p_rt.font.color.rgb = COLOR_PRIMARY
    p_rt.font.name = "맑은 고딕"
    
    right_content_box = slide2.shapes.add_textbox(Inches(7.2), Inches(2.3), Inches(5.2), Inches(4.2))
    tf_rc = right_content_box.text_frame
    tf_rc.word_wrap = True
    
    bullets_rc = [
        "정량/정성 데이터의 하이브리드 수집 파이프라인",
        "  - 정형 데이터(재무제표 등)는 DART API를 통해 오차율 0%로 추출",
        "  - 비정형 데이터(뉴스 등)는 실시간 뉴스 크롤러 및 임베딩 구성",
        "Multi-Agent 및 RAG 기반 정교화된 심사",
        "  - 수집 데이터 기반 부정 신호 센싱 및 실시간 리스크 게이지 산출",
        "  - RAG 결합으로 모든 리스크 문장에 대한 명확한 출처/원문 인용 제공",
        "안전한 금융 특화 하이브리드 LLM 가이드라인",
        "  - 대형 상용 LLM의 리서치 능력과 사내 On-Premise SLM 결합",
        "  - 심사역이 판단에만 고도로 집중하는 '심사 Copilot'으로 기능 설계"
    ]
    for i, b in enumerate(bullets_rc):
        p = tf_rc.paragraphs[0] if i == 0 else tf_rc.add_paragraph()
        p.text = b
        p.font.size = Pt(13) if b.startswith("  -") else Pt(14)
        p.font.bold = not b.startswith("  -")
        p.font.color.rgb = COLOR_TEXT_DARK
        p.font.name = "맑은 고딕"
        if b.startswith("  -"):
            p.font.color.rgb = COLOR_TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 3: 해결하고자 하는 현업/업무 과제 및 기대 효과
    # ----------------------------------------------------
    slide3 = add_slide_header("현업 과제 해결 및 비즈니스 기대 효과 (Expected Benefits)", "02")
    
    # 프로세스 전환 도표화 카드 (중앙 상단)
    process_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(12.1), Inches(1.8))
    process_box.fill.solid()
    process_box.fill.fore_color.rgb = COLOR_CARD_GRAY
    process_box.line.color.rgb = COLOR_PRIMARY
    process_box.line.width = Pt(1.0)
    
    proc_text_box = slide3.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(11.5), Inches(1.6))
    tf_proc = proc_text_box.text_frame
    tf_proc.word_wrap = True
    p_p1 = tf_proc.paragraphs[0]
    p_p1.text = "🔄 프로세스 패러다임 전환 (Paradigm Shift)"
    p_p1.font.size = Pt(16)
    p_p1.font.bold = True
    p_p1.font.color.rgb = COLOR_PRIMARY
    p_p1.font.name = "맑은 고딕"
    
    p_p2 = tf_proc.add_paragraph()
    p_p2.text = "  [기존 업무]  수작업 수집 ──> 데이터 일일이 수작업 타이핑 ──> 보고서 서식 수작업 작성 ➔ [최소 2~3일 소요]\n  [개선 업무]  대상 기업명 입력 ──> Multi-Agent 자동 수집/센싱 ──> Credit Memo 초안 10분 내 자동 도출 ➔ [심사역은 승인 판단만 수행]"
    p_p2.font.size = Pt(12)
    p_p2.font.bold = True
    p_p2.font.color.rgb = COLOR_TEXT_DARK
    p_p2.font.name = "맑은 고딕"
    p_p2.space_before = Pt(8)
    
    # 세부 업무 효과 및 기대 효과 (좌/우 배치)
    col1_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(3.6), Inches(5.8), Inches(3.2))
    col1_card.fill.solid()
    col1_card.fill.fore_color.rgb = COLOR_WHITE
    col1_card.line.color.rgb = COLOR_PRIMARY
    col1_card.line.width = Pt(1.0)
    
    col1_tb = slide3.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(5.4), Inches(2.8))
    tf_col1 = col1_tb.text_frame
    tf_col1.word_wrap = True
    p_c1 = tf_col1.paragraphs[0]
    p_c1.text = "📌 해결하고자 하는 현업 핵심 과제"
    p_c1.font.size = Pt(16)
    p_c1.font.bold = True
    p_c1.font.color.rgb = COLOR_PRIMARY
    p_c1.font.name = "맑은 고딕"
    
    bullets_c1 = [
        "• 정보 탐색 및 문서 타이핑 최소화: 단순 기입 자동화",
        "• 비정형 리스크 점량화: 기업 부정 뉴스 조기 감지 체계 구현",
        "• 부서 간 분석 격차 극복: 표준화된 심사 양식 자동 매핑"
    ]
    for b in bullets_c1:
        p = tf_col1.add_paragraph()
        p.text = b
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_DARK
        p.font.name = "맑은 고딕"
        p.space_before = Pt(6)
        
    col2_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(3.6), Inches(5.8), Inches(3.2))
    col2_card.fill.solid()
    col2_card.fill.fore_color.rgb = COLOR_WHITE
    col2_card.line.color.rgb = COLOR_PRIMARY
    col2_card.line.width = Pt(1.0)
    
    col2_tb = slide3.shapes.add_textbox(Inches(7.1), Inches(3.8), Inches(5.4), Inches(2.8))
    tf_col2 = col2_tb.text_frame
    tf_col2.word_wrap = True
    p_c2 = tf_col2.paragraphs[0]
    p_c2.text = "📈 기대 효과 및 재무 가치"
    p_c2.font.size = Pt(16)
    p_c2.font.bold = True
    p_c2.font.color.rgb = COLOR_PRIMARY
    p_c2.font.name = "맑은 고딕"
    
    bullets_c2 = [
        "• 생산성 약 90% 이상 개선: 심사 소요시간 획기적 단축",
        "• 부실여신 선제 방어: NPL(부실채권) 비율의 적극적 감소 유도",
        "• 법적 신뢰성 제고: AI에 원본 출처 인용 연계로 신뢰성 확보",
        "• 심사 전문성 상향 평준화: 누구나 일치하는 고품질 보고서 수립"
    ]
    for b in bullets_c2:
        p = tf_col2.add_paragraph()
        p.text = b
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_DARK
        p.font.name = "맑은 고딕"
        p.space_before = Pt(6)

    # ----------------------------------------------------
    # SLIDE 4: 사용 데이터 및 주요 기술 아키텍처
    # ----------------------------------------------------
    slide4 = add_slide_header("사용 예정 데이터 및 모델 아키텍처 (Data & Architecture)", "03")
    
    # 좌측 데이터 구성 카드
    data_card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3))
    data_card.fill.solid()
    data_card.fill.fore_color.rgb = COLOR_CARD_GRAY
    data_card.line.color.rgb = COLOR_PRIMARY
    data_card.line.width = Pt(1.0)
    
    data_tb = slide4.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(5.4), Inches(4.8))
    tf_data = data_tb.text_frame
    tf_data.word_wrap = True
    
    p_d = tf_data.paragraphs[0]
    p_d.text = "📊 사용 예정 데이터 수집 & 활용 계획"
    p_d.font.size = Pt(18)
    p_d.font.bold = True
    p_d.font.color.rgb = COLOR_PRIMARY
    p_d.font.name = "맑은 고딕"
    
    bullets_d = [
        "정형 데이터 (Hard Data): 환각 우려 원천 제거",
        "  - 대상: DART 공시 재무제표, 주요 재무비율, KRX 대주주 현황",
        "  - 획득: Open API 및 가상의 마스킹 데이터 활용",
        "  - 처리: LLM 배제, 오로지 규칙 기반 코딩으로 100% 정합성 수집",
        "비정형 데이터 (Soft Data): RAG 기반 분석",
        "  - 대상: 포털 뉴스 API, 우리금융경영연구소 산업 보고서",
        "  - 획득: BeautifulSoup/Selenium 및 공용 데이터포털 연계",
        "  - 처리: 텍스트 임베딩 후 Chroma DB(Vector DB) 적재",
        "금융 보안 가이드라인 준수",
        "  - 가상 기업 정보(10개사) 기반으로 가명 처리하여 유출 위험 제로"
    ]
    for b in bullets_d:
        p = tf_data.add_paragraph()
        p.text = b
        p.font.size = Pt(12) if b.startswith("  -") else Pt(13)
        p.font.bold = not b.startswith("  -")
        p.font.color.rgb = COLOR_TEXT_DARK
        p.font.name = "맑은 고딕"
        if b.startswith("  -"):
            p.font.color.rgb = COLOR_TEXT_MUTED
            
    # 우측 아키텍처 구성 카드
    arch_card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.3))
    arch_card.fill.solid()
    arch_card.fill.fore_color.rgb = COLOR_CARD_GRAY
    arch_card.line.color.rgb = COLOR_PRIMARY
    arch_card.line.width = Pt(1.0)
    
    arch_tb = slide4.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.4), Inches(4.8))
    tf_arch = arch_tb.text_frame
    tf_arch.word_wrap = True
    
    p_a = tf_arch.paragraphs[0]
    p_a.text = "⚙️ 시스템 기술 아키텍처 및 도구"
    p_a.font.size = Pt(18)
    p_a.font.bold = True
    p_a.font.color.rgb = COLOR_PRIMARY
    p_a.font.name = "맑은 고딕"
    
    bullets_a = [
        "Multi-Agent 제어 (LangGraph): 흐름 정교화",
        "  - Data Agent: DART API를 연동하여 재무 수치 분석",
        "  - Risk Agent: 부정 뉴스 감성 분석 및 리스크 스코어 계산",
        "  - Drafting Agent: 결과 결합 후 Credit Memo 최종 조립",
        "Hybrid LLM 운영 정책: 보안과 성능의 합의점",
        "  - 외부망 GPT-4o: 비정형 데이터 정제 및 추론 성능 극대화",
        "  - On-Premise SLM: 사내 내부망 템플릿 문서 빌드용으로 활용",
        "시연용 인터페이스 (Streamlit)",
        "  - 기업명 입력 대시보드, 리스크 게이지 시각화, 보고서 다운로드"
    ]
    for b in bullets_a:
        p = tf_arch.add_paragraph()
        p.text = b
        p.font.size = Pt(12) if b.startswith("  -") else Pt(13)
        p.font.bold = not b.startswith("  -")
        p.font.color.rgb = COLOR_TEXT_DARK
        p.font.name = "맑은 고딕"
        if b.startswith("  -"):
            p.font.color.rgb = COLOR_TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 5: 추진 일정 및 평가지표
    # ----------------------------------------------------
    slide5 = add_slide_header("개발 및 추진 계획과 핵심 평가 지표 (Roadmap & KPI)", "04")
    
    # 좌측 로드맵 카드 (일정)
    sched_card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3))
    sched_card.fill.solid()
    sched_card.fill.fore_color.rgb = COLOR_WHITE
    sched_card.line.color.rgb = COLOR_PRIMARY
    sched_card.line.width = Pt(1.0)
    
    sched_tb = slide5.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(5.4), Inches(4.8))
    tf_sched = sched_tb.text_frame
    tf_sched.word_wrap = True
    
    p_s = tf_sched.paragraphs[0]
    p_s.text = "📅 12주 현실적 개발 로드맵 (6~18주차)"
    p_s.font.size = Pt(18)
    p_s.font.bold = True
    p_s.font.color.rgb = COLOR_PRIMARY
    p_s.font.name = "맑은 고딕"
    
    milestones = [
        "6 ~ 8주차: 데이터 파이프라인 수립 및 데이터 가상화",
        "  - DART, 뉴스 API 연동 모듈 및 가상 10개사 데이터셋 구축",
        "9 ~ 11주차: Multi-Agent 및 RAG 아키텍처 모델링",
        "  - LangGraph 에이전트 설계, 뉴스 분석 RAG 및 VectorDB 이식",
        "12 ~ 14주차: 보고서 파일 빌더 및 Streamlit 대시보드 연동",
        "  - Markdown/PDF 파일 작성 로직 정합성 점검 및 UI 고도화",
        "15 ~ 17주차: End-to-End 통합 검증 및 프롬프트 튜닝",
        "  - 에이전트 연동 에러 방지 및 환각(Hallucination) 제어 수치 튜닝",
        "18주차: 최종 데모 영상 및 발표 자료 최종화"
    ]
    for b in milestones:
        p = tf_sched.add_paragraph()
        p.text = b
        p.font.size = Pt(11) if b.startswith("  -") else Pt(13)
        p.font.bold = not b.startswith("  -")
        p.font.color.rgb = COLOR_TEXT_DARK
        p.font.name = "맑은 고딕"
        if b.startswith("  -"):
            p.font.color.rgb = COLOR_TEXT_MUTED
            
    # 우측 평가 지표 카드 (KPI)
    kpi_card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.3))
    kpi_card.fill.solid()
    kpi_card.fill.fore_color.rgb = COLOR_CARD_BLUE
    kpi_card.line.color.rgb = COLOR_PRIMARY
    kpi_card.line.width = Pt(1.5)
    
    kpi_tb = slide5.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.4), Inches(4.8))
    tf_kpi = kpi_tb.text_frame
    tf_kpi.word_wrap = True
    
    p_k = tf_kpi.paragraphs[0]
    p_k.text = "🎯 핵심 평가지표 및 검증 타당성"
    p_k.font.size = Pt(18)
    p_k.font.bold = True
    p_k.font.color.rgb = COLOR_PRIMARY
    p_k.font.name = "맑은 고딕"
    
    bullets_k = [
        "정량적 평가 지표 (Technical Performance)",
        "  1. 재무 정보 정합성 (Data Accuracy): DART 수치 일치율 100%",
        "  2. 보고서 도출 속도: 수동 2~3일 ➔ 자동 10분 이내 (90% 이상 향상)",
        "  3. 생성 도달율: 파이프라인 에러율 5% 미만 (성공률 95% 이상)",
        "정성적 평가 지표 (Expert Verification)",
        "  1. 문장 원문 근거성 (BERTScore): 소스 연동 평가 0.82 이상",
        "  2. 현업 블라인드 피델리티 평가 (Fidelity Test):",
        "       - 우리은행 여신 전문가 3인의 보고서 수준 만족도 블라인드 테스트",
        "       - 목표 평점 4.0점 이상 / 5.0점 만점 획득"
    ]
    for b in bullets_k:
        p = tf_kpi.add_paragraph()
        p.text = b
        p.font.size = Pt(11) if (b.startswith("  1") or b.startswith("  2") or b.startswith("  3") or b.startswith("       -")) else Pt(13)
        p.font.bold = not (b.startswith("  ") or b.startswith("       -"))
        p.font.color.rgb = COLOR_TEXT_DARK
        p.font.name = "맑은 고딕"
        if b.startswith("  1") or b.startswith("  2") or b.startswith("  3") or b.startswith("       -"):
            p.font.color.rgb = COLOR_TEXT_MUTED

    # 발표 끝 (Thank you) 슬라이드 생략하고 저장
    prs.save("우리금융_개인프로젝트_제안발표_슬라이드.pptx")
    print("SUCCESS: PPTX presentation successfully generated.")

if __name__ == "__main__":
    create_presentation()
